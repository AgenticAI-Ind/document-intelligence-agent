"""
Multi-format document parsing (PDF, Word, Excel, PowerPoint)
"""

import logging
from typing import Dict, List, Optional, BinaryIO
from pathlib import Path
import io

logger = logging.getLogger(__name__)


class DocumentParser:
    """Parse documents from multiple formats"""

    def __init__(self):
        """Initialize document parser with support for multiple formats"""
        self.supported_formats = [
            'pdf', 'docx', 'doc', 'xlsx', 'xls',
            'pptx', 'ppt', 'txt', 'csv', 'md'
        ]

    def parse_document(
        self,
        file_path: Optional[str] = None,
        file_content: Optional[bytes] = None,
        file_extension: Optional[str] = None
    ) -> Dict:
        """
        Parse document from file path or bytes

        Args:
            file_path: Path to document file
            file_content: Raw bytes of document
            file_extension: File extension (required if using file_content)

        Returns:
            Dict with parsed content, metadata, and structure
        """
        try:
            # Determine file type
            if file_path:
                ext = Path(file_path).suffix.lower().lstrip('.')
            elif file_extension:
                ext = file_extension.lower().lstrip('.')
            else:
                raise ValueError("Must provide either file_path or file_extension")

            if ext not in self.supported_formats:
                raise ValueError(f"Unsupported format: {ext}")

            # Parse based on format
            if ext == 'pdf':
                return self._parse_pdf(file_path, file_content)
            elif ext in ['docx', 'doc']:
                return self._parse_word(file_path, file_content)
            elif ext in ['xlsx', 'xls']:
                return self._parse_excel(file_path, file_content)
            elif ext in ['pptx', 'ppt']:
                return self._parse_powerpoint(file_path, file_content)
            elif ext == 'txt':
                return self._parse_text(file_path, file_content)
            elif ext == 'csv':
                return self._parse_csv(file_path, file_content)
            elif ext == 'md':
                return self._parse_markdown(file_path, file_content)
            else:
                raise ValueError(f"Parser not implemented for: {ext}")

        except Exception as e:
            logger.error(f"Error parsing document: {e}")
            return {
                "content": "",
                "metadata": {},
                "error": str(e)
            }

    def _parse_pdf(
        self,
        file_path: Optional[str],
        file_content: Optional[bytes]
    ) -> Dict:
        """Parse PDF document"""
        try:
            import PyPDF2

            if file_path:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    return self._extract_pdf_content(reader)
            else:
                reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                return self._extract_pdf_content(reader)

        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            return {"content": "", "error": str(e)}

    def _extract_pdf_content(self, reader: 'PyPDF2.PdfReader') -> Dict:
        """Extract content from PDF reader"""
        pages = []
        full_text = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            pages.append({
                "page_number": i + 1,
                "text": text,
                "char_count": len(text)
            })
            full_text.append(text)

        metadata = reader.metadata or {}

        return {
            "content": "\n\n".join(full_text),
            "pages": pages,
            "page_count": len(pages),
            "metadata": {
                "title": metadata.get('/Title', ''),
                "author": metadata.get('/Author', ''),
                "subject": metadata.get('/Subject', ''),
                "creator": metadata.get('/Creator', ''),
                "producer": metadata.get('/Producer', ''),
                "creation_date": str(metadata.get('/CreationDate', '')),
            },
            "format": "pdf"
        }

    def _parse_word(
        self,
        file_path: Optional[str],
        file_content: Optional[bytes]
    ) -> Dict:
        """Parse Word document (.docx)"""
        try:
            from docx import Document

            if file_path:
                doc = Document(file_path)
            else:
                doc = Document(io.BytesIO(file_content))

            # Extract paragraphs
            paragraphs = []
            full_text = []

            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append({
                        "text": para.text,
                        "style": para.style.name
                    })
                    full_text.append(para.text)

            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)

            # Core properties
            core_props = doc.core_properties

            return {
                "content": "\n\n".join(full_text),
                "paragraphs": paragraphs,
                "tables": tables,
                "table_count": len(tables),
                "paragraph_count": len(paragraphs),
                "metadata": {
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "created": str(core_props.created) if core_props.created else "",
                    "modified": str(core_props.modified) if core_props.modified else "",
                },
                "format": "docx"
            }

        except Exception as e:
            logger.error(f"Error parsing Word document: {e}")
            return {"content": "", "error": str(e)}

    def _parse_excel(
        self,
        file_path: Optional[str],
        file_content: Optional[bytes]
    ) -> Dict:
        """Parse Excel spreadsheet"""
        try:
            import openpyxl

            if file_path:
                wb = openpyxl.load_workbook(file_path, data_only=True)
            else:
                wb = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)

            sheets = []
            all_text = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]

                # Extract data
                data = []
                for row in ws.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):  # Skip empty rows
                        data.append(row_data)
                        all_text.append(" ".join(row_data))

                sheets.append({
                    "name": sheet_name,
                    "data": data,
                    "row_count": len(data),
                    "col_count": len(data[0]) if data else 0
                })

            return {
                "content": "\n".join(all_text),
                "sheets": sheets,
                "sheet_count": len(sheets),
                "metadata": {
                    "title": wb.properties.title or "",
                    "author": wb.properties.creator or "",
                    "created": str(wb.properties.created) if wb.properties.created else "",
                },
                "format": "xlsx"
            }

        except Exception as e:
            logger.error(f"Error parsing Excel: {e}")
            return {"content": "", "error": str(e)}

    def _parse_powerpoint(
        self,
        file_path: Optional[str],
        file_content: Optional[bytes]
    ) -> Dict:
        """Parse PowerPoint presentation"""
        try:
            from pptx import Presentation

            if file_path:
                prs = Presentation(file_path)
            else:
                prs = Presentation(io.BytesIO(file_content))

            slides = []
            all_text = []

            for i, slide in enumerate(prs.slides):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                            all_text.append(text)

                slides.append({
                    "slide_number": i + 1,
                    "text": "\n".join(slide_text),
                    "shape_count": len(slide.shapes)
                })

            return {
                "content": "\n\n".join(all_text),
                "slides": slides,
                "slide_count": len(slides),
                "metadata": {
                    "title": prs.core_properties.title or "",
                    "author": prs.core_properties.author or "",
                    "subject": prs.core_properties.subject or "",
                },
                "format": "pptx"
            }

        except Exception as e:
            logger.error(f"Error parsing PowerPoint: {e}")
            return {"content": "", "error": str(e)}

    def _parse_text(
        self,
        file_path: Optional[str],
        file_content: Optional[bytes]
    ) -> Dict:
        """Parse plain text file"""
        try:
            if file_path:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            else:
                content = file_content.decode('utf-8')

            lines = content.split('\n')

            return {
                "content": content,
                "line_count": len(lines),
                "char_count": len(content),
                "word_count": len(content.split()),
                "format": "txt"
            }

        except Exception as e:
            logger.error(f"Error parsing text: {e}")
            return {"content": "", "error": str(e)}

    def _parse_csv(
        self,
        file_path: Optional[str],
        file_content: Optional[bytes]
    ) -> Dict:
        """Parse CSV file"""
        try:
            import csv

            if file_path:
                with open(file_path, 'r', encoding='utf-8') as f:
                    reader = csv.reader(f)
                    rows = list(reader)
            else:
                content = file_content.decode('utf-8')
                reader = csv.reader(io.StringIO(content))
                rows = list(reader)

            # Convert to text
            text_rows = [", ".join(row) for row in rows]

            return {
                "content": "\n".join(text_rows),
                "rows": rows,
                "row_count": len(rows),
                "col_count": len(rows[0]) if rows else 0,
                "format": "csv"
            }

        except Exception as e:
            logger.error(f"Error parsing CSV: {e}")
            return {"content": "", "error": str(e)}

    def _parse_markdown(
        self,
        file_path: Optional[str],
        file_content: Optional[bytes]
    ) -> Dict:
        """Parse Markdown file"""
        try:
            if file_path:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            else:
                content = file_content.decode('utf-8')

            # Extract headers
            import re
            headers = re.findall(r'^#{1,6}\s+(.+)$', content, re.MULTILINE)

            return {
                "content": content,
                "headers": headers,
                "header_count": len(headers),
                "char_count": len(content),
                "format": "md"
            }

        except Exception as e:
            logger.error(f"Error parsing markdown: {e}")
            return {"content": "", "error": str(e)}

    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return self.supported_formats.copy()

    def is_format_supported(self, extension: str) -> bool:
        """Check if a file format is supported"""
        return extension.lower().lstrip('.') in self.supported_formats
